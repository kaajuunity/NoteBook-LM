import os
import json
import asyncio
import uuid
import glob
import secrets
import psycopg2
from psycopg2.extras import RealDictCursor
from flask import Flask, request, jsonify, render_template, session

import google.generativeai as genai
import edge_tts
from pydub import AudioSegment
import pypdf

from config import Config
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
from PIL import Image, ImageDraw, ImageFont
import textwrap
import shutil


app = Flask(__name__)
app.secret_key = os.getenv('SECRET_KEY', secrets.token_hex(32))

# Configuration
genai.configure(api_key=Config.GEMINI_API_KEY)
model = genai.GenerativeModel(
    model_name=Config.MODEL_NAME,
    generation_config=Config.GENERATION_CONFIG,
)

def get_db_connection():
    """Connect to Supabase PostgreSQL with connection pooling"""
    conn = psycopg2.connect(
        Config.DATABASE_URL,
        sslmode='require',
        connect_timeout=10
    )
    return conn

def init_db():
    """Initialize database with required table and pgvector extension"""
    conn = None
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        

        cur.execute("CREATE EXTENSION IF NOT EXISTS vector")
        

        cur.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                id SERIAL PRIMARY KEY,
                content TEXT NOT NULL,
                metadata JSONB,
                embedding vector(768),
                user_id TEXT NOT NULL DEFAULT 'default_user',
                project_id TEXT NOT NULL DEFAULT 'default_project',
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        

        cur.execute("""
            CREATE INDEX IF NOT EXISTS documents_embedding_idx 
            ON documents USING ivfflat (embedding vector_cosine_ops)
            WITH (lists = 100)
        """)
        
        cur.execute("""
            CREATE INDEX IF NOT EXISTS idx_user_project 
            ON documents(user_id, project_id)
        """)
        
        cur.execute("""
            CREATE INDEX IF NOT EXISTS idx_created_at 
            ON documents(created_at DESC)
        """)
        
        conn.commit()
        cur.close()
        print("✓ Database initialized successfully on Supabase")
    except Exception as e:
        print(f"Database initialization: {e}")
        if conn:
            conn.rollback()
    finally:
        if conn:
            conn.close()

def get_user_session():
    """Get or create user session with project ID"""
    if 'user_id' not in session:
        session['user_id'] = f"user_{uuid.uuid4().hex[:12]}"
    if 'project_id' not in session:
        session['project_id'] = f"project_{uuid.uuid4().hex[:12]}"
    return session['user_id'], session['project_id']

# Helpers
def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunks.append(text[start:end])
        start += (chunk_size - overlap)
    return chunks

def generate_audio_with_gtts(text, output_file):
    """Fallback audio generation using Google TTS"""
    try:
        from gtts import gTTS
        tts = gTTS(text=text, lang='en', slow=False)
        tts.save(output_file)
        return True
    except Exception as e:
        print(f"gTTS error: {e}")
        return False

async def generate_audio_clip(text, voice, output_file):
    try:
        communicate = edge_tts.Communicate(text, voice)
        await communicate.save(output_file)
        return True
    except Exception as e:
        print(f"Audio generation error for {output_file}: {e}")
        return False

async def generate_all_audio_clips(script_json, temp_dir):
    """Generate all audio clips using proper async approach"""
    audio_files = []
    
    for i, turn in enumerate(script_json):
        speaker = turn.get('speaker', 'Host A')
        text = turn.get('text', '')
        if not text:
            continue
            
        voice = "en-US-GuyNeural" if speaker == "Host A" else "en-US-JennyNeural"
        filename = f"{temp_dir}/chunk_{i}.mp3"
        
        success = await generate_audio_clip(text, voice, filename)
        if success and os.path.exists(filename):
            audio_files.append(filename)
    
    return audio_files

# Routes
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    conn = None
    try:
        file = request.files.get('file')
        if not file:
            return jsonify({'error': 'No file uploaded'}), 400
        
        filename = file.filename
        text_content = ""
        
        if filename.endswith('.pdf'):
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content += page_text
        else:

            text_content = file.read().decode('utf-8')

        if not text_content.strip():
            return jsonify({'error': 'No text content found in file'}), 400

        chunks = chunk_text(text_content)
        
        if len(chunks) > 500:
            return jsonify({'error': 'File too large. Please upload smaller files.'}), 400
        

        user_id, project_id = get_user_session()
        
        conn = get_db_connection()
        cur = conn.cursor()

        processed_count = 0
        for chunk in chunks:

            embedding_result = genai.embed_content(
                model=Config.EMBEDDING_MODEL,
                content=chunk,
                task_type="retrieval_document"
            )
            embedding = embedding_result['embedding']
            

            if len(embedding) != 768:
                raise ValueError(f"Unexpected embedding dimension: {len(embedding)}")
            

            cur.execute(
                "INSERT INTO documents (content, metadata, embedding, user_id, project_id) VALUES (%s, %s, %s, %s, %s)",
                (chunk, json.dumps({'filename': filename}), embedding, user_id, project_id)
            )
            processed_count += 1
            
        conn.commit()
        cur.close()
        
        print(f"✓ Uploaded {processed_count} chunks to Supabase for user: {user_id}, project: {project_id}")
        
        return jsonify({'message': f'Successfully processed {processed_count} chunks from {filename}'}), 200

    except Exception as e:
        print(f"Upload error: {e}")
        if conn:
            conn.rollback()
        return jsonify({'error': 'Failed to process file. Please try again.'}), 500
    finally:
        if conn:
            conn.close()

@app.route('/chat', methods=['POST'])
def chat():
    conn = None
    try:
        data = request.json
        user_query = data.get('query')
        if not user_query:
            return jsonify({'error': 'Query is required'}), 400


        user_id, project_id = get_user_session()


        embedding_result = genai.embed_content(
            model=Config.EMBEDDING_MODEL,
            content=user_query,
            task_type="retrieval_query"
        )
        query_embedding = embedding_result['embedding']

        print(f"Query embedding size: {len(query_embedding)}")

        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)


        cur.execute("""
            SELECT content, 1 - (embedding <=> %s::vector) as similarity
            FROM documents
            WHERE user_id = %s 
            AND project_id = %s
            AND 1 - (embedding <=> %s::vector) > 0.3
            ORDER BY similarity DESC
            LIMIT 5
        """, (query_embedding, user_id, project_id, query_embedding))
        
        results = cur.fetchall()
        
        print(f"Found {len(results)} matching documents for user: {user_id}, project: {project_id}")
        if results:
            print(f"Best similarity score: {results[0]['similarity']}")
        
        cur.close()
        
        if not results:
            return jsonify({'answer': 'I couldn\'t find relevant information in your sources to answer this question.'})
        
        relevant_chunks = [row['content'] for row in results]
        context = "\n\n".join(relevant_chunks)
        
        prompt = f"""You are a helpful assistant. Answer the user question strictly based on the provided context.
        
Context:
{context}

Question: {user_query}

Answer:"""
        
        response = model.generate_content(prompt)
        return jsonify({'answer': response.text})

    except Exception as e:
        print(f"Chat error: {e}")
        error_msg = str(e)
        if "429" in error_msg or "quota" in error_msg.lower():
            return jsonify({'error': 'API quota exceeded. Please wait a moment and try again.'}), 429
        return jsonify({'error': 'Failed to process your question. Please try again.'}), 500
    finally:
        if conn:
            conn.close()

@app.route('/generate_audio', methods=['POST'])
def generate_audio():
    conn = None
    try:

        user_id, project_id = get_user_session()
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        

        cur.execute("""
            SELECT content FROM documents 
            WHERE user_id = %s AND project_id = %s
            ORDER BY created_at DESC 
            LIMIT 50
        """, (user_id, project_id))
        rows = cur.fetchall()
        cur.close()
        
        if not rows:
            return jsonify({'error': 'No content available to generate audio'}), 400
        
        all_text = " ".join([row['content'] for row in rows])
        

        if len(all_text) > 50000:
            all_text = all_text[:50000]
            last_period = all_text.rfind('.')
            if last_period > 40000:
                all_text = all_text[:last_period + 1]

        script_prompt = f"""Generate a podcast script between two hosts (Host A and Host B) discussing this content. 
Make it conversational, engaging, and use simple English. Keep it concise (max 8-10 exchanges).
Format as JSON: [{{"speaker": "Host A", "text": "..."}}, {{"speaker": "Host B", "text": "..."}}].

Content:
{all_text}
"""
        
        model_config = {"response_mime_type": "application/json"}
        script_model = genai.GenerativeModel(Config.MODEL_NAME, generation_config=model_config)
        
        script_response = script_model.generate_content(script_prompt)
        script_json = json.loads(script_response.text)
        
        temp_dir = "temp_audio"
        os.makedirs(temp_dir, exist_ok=True)
        

        audio_files = asyncio.run(generate_all_audio_clips(script_json, temp_dir))
        
        if not audio_files:
            return jsonify({'error': 'Failed to generate audio clips'}), 500
        

        combined_audio = AudioSegment.empty()
        for filename in audio_files:
            try:
                segment = AudioSegment.from_mp3(filename)
                combined_audio += segment
            except Exception as e:
                print(f"Error loading audio segment {filename}: {e}")
        
        output_filename = f"audio_overview_{uuid.uuid4()}.mp3"
        output_path = os.path.join("static", output_filename)
        combined_audio.export(output_path, format="mp3")
        

        for f in audio_files:
            try:
                os.remove(f)
            except:
                pass
        
        return jsonify({'audio_url': f'/static/{output_filename}'})

    except Exception as e:
        print(f"Error in generate_audio: {e}")
        return jsonify({'error': 'Failed to generate audio. Please try again.'}), 500
    finally:
        if conn:
            conn.close()

@app.route('/generate_study_aid', methods=['POST'])
def generate_study_aid():
    conn = None
    try:
        data = request.json
        aid_type = data.get('type')
        
        if aid_type not in ['flowchart', 'flashcard', 'quiz']:
            return jsonify({'error': 'Invalid type'}), 400
        

        user_id, project_id = get_user_session()
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        

        cur.execute("""
            SELECT content FROM documents 
            WHERE user_id = %s AND project_id = %s
            ORDER BY created_at DESC 
            LIMIT 50
        """, (user_id, project_id))
        rows = cur.fetchall()
        cur.close()
        
        if not rows:
            return jsonify({'error': 'No content available'}), 400
        
        all_text = " ".join([row['content'] for row in rows])
        

        if len(all_text) > 50000:
            all_text = all_text[:50000]
            last_period = all_text.rfind('.')
            if last_period > 40000:
                all_text = all_text[:last_period + 1]
        
        if aid_type == 'flowchart':
            prompt = f"""Generate Mermaid.js code representing the key concepts and their relationships in this text. 
Return ONLY the mermaid code starting with 'graph TD' or 'graph LR'. Do not include markdown code fences or any other text.

Text: {all_text}"""
            
            resp = model.generate_content(prompt)
            clean_text = resp.text.replace('```mermaid', '').replace('```', '').strip()
            return jsonify({'content': clean_text})
            
        elif aid_type in ['flashcard', 'quiz']:
            if aid_type == 'quiz':
                prompt = f"""Generate a JSON array of 5-7 multiple choice questions based on the text.
Each question should have:
- question: the question text
- options: array of 4 possible answers
- answer: the correct answer (must be one of the options)
- explanation: brief explanation of why this is correct (1-2 sentences)
- wrongExplanation: brief explanation shown for wrong answers (1-2 sentences)

Format: [{{"question": "...", "options": ["A", "B", "C", "D"], "answer": "A", "explanation": "...", "wrongExplanation": "..."}}]
Make questions clear and test understanding of key concepts.

Text: {all_text}"""
            else:
                prompt = f"""Generate a JSON array of 5-10 questions and answers for a {aid_type} based on the text.
Format: [{{"question": "...", "answer": "..."}}]
Make questions clear and answers concise.

Text: {all_text}"""
        
        model_config = {"response_mime_type": "application/json"}
        json_model = genai.GenerativeModel(Config.MODEL_NAME, generation_config=model_config)
        resp = json_model.generate_content(prompt)
        return jsonify(json.loads(resp.text))

    except Exception as e:
        print(f"Study aid error: {e}")
        return jsonify({'error': 'Failed to generate study aid. Please try again.'}), 500
    finally:
        if conn:
            conn.close()

@app.route('/generate_slides', methods=['POST'])
def generate_slides():
    conn = None
    try:

        user_id, project_id = get_user_session()
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        

        cur.execute("""
            SELECT content FROM documents 
            WHERE user_id = %s AND project_id = %s
            ORDER BY created_at DESC 
            LIMIT 50
        """, (user_id, project_id))
        rows = cur.fetchall()
        cur.close()
        
        if not rows:
            return jsonify({'error': 'No content available'}), 400
        
        all_text = " ".join([row['content'] for row in rows])
        

        if len(all_text) > 50000:
            all_text = all_text[:50000]
            last_period = all_text.rfind('.')
            if last_period > 40000:
                all_text = all_text[:last_period + 1]
        

        prompt = f"""Generate a professional presentation with 6-8 slides based on this content.
Return ONLY valid JSON in this exact format:
{{
    "title": "Main Presentation Title",
    "slides": [
        {{
            "type": "title",
            "title": "Main Title",
            "subtitle": "Subtitle text"
        }},
        {{
            "type": "content",
            "title": "Slide Title",
            "points": ["Point 1", "Point 2", "Point 3", "Point 4"]
        }}
    ]
}}

Rules:
- First slide must be type "title" with title and subtitle
- Other slides must be type "content" with title and 3-5 points
- Each point should be concise (max 15 words)
- Make it professional and well-structured
- Focus on key concepts and important information

Content:
{all_text}"""
        
        model_config = {"response_mime_type": "application/json"}
        json_model = genai.GenerativeModel(Config.MODEL_NAME, generation_config=model_config)
        resp = json_model.generate_content(prompt)
        slide_data = json.loads(resp.text)
        

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for slide_info in slide_data.get('slides', []):
            if slide_info['type'] == 'title':

                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = slide_info.get('title', 'Untitled')
                subtitle.text = slide_info.get('subtitle', '')
                
            elif slide_info['type'] == 'content':

                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = slide_info.get('title', 'Slide')
                

                body_shape = slide.shapes.placeholders[1]
                text_frame = body_shape.text_frame
                text_frame.clear()
                
                for point in slide_info.get('points', []):
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(18)
        

        ppt_filename = f"presentation_{uuid.uuid4()}.pptx"
        ppt_path = os.path.join("static", ppt_filename)
        prs.save(ppt_path)
        

        return jsonify({
            'slides': slide_data,
            'download_url': f'/static/{ppt_filename}'
        })

    except Exception as e:
        print(f"Slide generation error: {e}")
        return jsonify({'error': 'Failed to generate slides. Please try again.'}), 500
    finally:
        if conn:
            conn.close()

@app.route('/generate_video', methods=['POST'])
def generate_video():
    conn = None
    try:

        user_id, project_id = get_user_session()
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        

        cur.execute("""
            SELECT content FROM documents 
            WHERE user_id = %s AND project_id = %s
            ORDER BY created_at DESC 
            LIMIT 50
        """, (user_id, project_id))
        rows = cur.fetchall()
        cur.close()
        
        if not rows:
            return jsonify({'error': 'No content available'}), 400
        
        all_text = " ".join([row['content'] for row in rows])
        

        if len(all_text) > 50000:
            all_text = all_text[:50000]
            last_period = all_text.rfind('.')
            if last_period > 40000:
                all_text = all_text[:last_period + 1]
        

        video_prompt = f"""Generate a video presentation with 6-8 slides and matching narration.
Target duration: 3-5 minutes total.

Return ONLY valid JSON in this exact format:
{{
    "title": "Main Presentation Title",
    "slides": [
        {{
            "type": "title",
            "title": "Main Title",
            "subtitle": "Subtitle text",
            "narration": "Welcome to this presentation about..."
        }},
        {{
            "type": "content",
            "title": "Slide Title",
            "points": ["Point 1", "Point 2", "Point 3"],
            "narration": "In this section, we'll explore..."
        }}
    ]
}}

Rules:
- First slide must be type "title" with title, subtitle, and narration
- Other slides must be type "content" with title, 3-5 points, and narration
- Each narration should be 15-25 seconds of speech (about 40-65 words)
- Narration should flow naturally and explain the slide content
- Keep points concise (max 12 words each)
- Make it professional and engaging

Content:
{all_text}"""
        
        model_config = {"response_mime_type": "application/json"}
        json_model = genai.GenerativeModel(Config.MODEL_NAME, generation_config=model_config)
        resp = json_model.generate_content(video_prompt)
        video_data = json.loads(resp.text)
        

        temp_dir = "temp_video"
        slide_dir = os.path.join(temp_dir, "slides")
        audio_dir = os.path.join(temp_dir, "audio")
        os.makedirs(slide_dir, exist_ok=True)
        os.makedirs(audio_dir, exist_ok=True)
        

        slide_images = []
        for i, slide_info in enumerate(video_data.get('slides', [])):
            img_path = os.path.join(slide_dir, f"slide_{i}.png")
            create_slide_image(slide_info, img_path)
            slide_images.append(img_path)
        

        audio_files = []
        for i, slide_info in enumerate(video_data.get('slides', [])):
            narration = slide_info.get('narration', '').strip()
            

            if not narration or len(narration) < 10:
                print(f"⚠ Slide {i}: No narration text")
                audio_files.append(None)
                continue
            

            narration = narration.replace('"', '').replace("'", "").replace('\n', ' ')
            narration = ' '.join(narration.split())  
            
            audio_path = os.path.join(audio_dir, f"narration_{i}.mp3")
            audio_success = False
            

            try:
                print(f"Trying Edge TTS for slide {i}...")
                success = asyncio.run(generate_audio_clip(narration, "en-US-GuyNeural", audio_path))
                if success and os.path.exists(audio_path) and os.path.getsize(audio_path) > 0:
                    audio_files.append(audio_path)
                    print(f"✓ Edge TTS success for slide {i}")
                    audio_success = True
            except Exception as e:
                print(f"⚠ Edge TTS failed for slide {i}: {str(e)[:100]}")
            

            if not audio_success:
                try:
                    print(f"Trying Google TTS for slide {i}...")
                    success = generate_audio_with_gtts(narration, audio_path)
                    if success and os.path.exists(audio_path) and os.path.getsize(audio_path) > 0:
                        audio_files.append(audio_path)
                        print(f"✓ Google TTS success for slide {i}")
                        audio_success = True
                except Exception as e:
                    print(f"⚠ Google TTS also failed for slide {i}: {str(e)[:100]}")
            

            if not audio_success:
                print(f"✗ All audio methods failed for slide {i}, using silent 10s duration")
                audio_files.append(None)
        

        clips = []
        for i, (img_path, audio_path) in enumerate(zip(slide_images, audio_files)):
            if audio_path and os.path.exists(audio_path):

                audio_clip = AudioFileClip(audio_path)
                duration = audio_clip.duration
                audio_clip.close()
            else:

                duration = 10
            

            img_clip = ImageClip(img_path, duration=duration)
            

            if audio_path and os.path.exists(audio_path):
                audio = AudioFileClip(audio_path)
                img_clip = img_clip.set_audio(audio)
            
            clips.append(img_clip)
        

        final_video = concatenate_videoclips(clips, method="compose")
        

        video_filename = f"video_overview_{uuid.uuid4()}.mp4"
        video_path = os.path.join("static", video_filename)
        
        final_video.write_videofile(
            video_path,
            fps=30, 
            codec='libx264',
            audio_codec='aac',
            preset='medium',  
            bitrate='8000k',  
            threads=4
        )
        

        video_duration = final_video.duration
        

        final_video.close()
        for clip in clips:
            clip.close()
        

        try:
            shutil.rmtree(temp_dir)
        except Exception as cleanup_error:
            print(f"Cleanup error: {cleanup_error}")
        
        return jsonify({
            'video_url': f'/static/{video_filename}',
            'duration': float(video_duration),
            'slides_count': len(slide_images)
        })
    
    except Exception as e:
        print(f"Video generation error: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': 'Failed to generate video. Please try again.'}), 500
    finally:
        if conn:
            conn.close()


def create_slide_image(slide_info, output_path, width=2560, height=1440):
    """Create a slide image from slide data"""

    img = Image.new('RGB', (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    

    try:
        title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 480)
        subtitle_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 360)
        body_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 340)
        point_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 220)
    except:

        title_font = ImageFont.load_default()
        subtitle_font = ImageFont.load_default()
        body_font = ImageFont.load_default()
        point_font = ImageFont.load_default()
    
    slide_type = slide_info.get('type', 'content')
    
    if slide_type == 'title':

        title = slide_info.get('title', 'Title')
        subtitle = slide_info.get('subtitle', '')
        

        title_bbox = draw.textbbox((0, 0), title, font=title_font)
        title_width = title_bbox[2] - title_bbox[0]
        title_height = title_bbox[3] - title_bbox[1]
        title_x = (width - title_width) // 2
        title_y = height // 2 - 100
        
        draw.text((title_x, title_y), title, fill=(0, 0, 0), font=title_font)
        

        if subtitle:
            subtitle_bbox = draw.textbbox((0, 0), subtitle, font=subtitle_font)
            subtitle_width = subtitle_bbox[2] - subtitle_bbox[0]
            subtitle_x = (width - subtitle_width) // 2
            subtitle_y = title_y + title_height + 50
            
            draw.text((subtitle_x, subtitle_y), subtitle, fill=(80, 80, 80), font=subtitle_font)
    
    else:
 
        title = slide_info.get('title', 'Slide')
        points = slide_info.get('points', [])
        
  
        title_y = 120
        draw.text((100, title_y), title, fill=(0, 0, 0), font=body_font)
        
 
        title_bbox = draw.textbbox((100, title_y), title, font=body_font)
        title_width = title_bbox[2] - title_bbox[0]
        draw.rectangle([(100, title_y + 70), (100 + title_width, title_y + 76)], fill=(0, 0, 0))
        

        y_position = title_y + 140
        for point in points:

            wrapped = textwrap.fill(point, width=110)
            lines = wrapped.split('\n')
            
            for line_idx, line in enumerate(lines):
                if line_idx == 0:

                    draw.ellipse([(120, y_position + 12), (145, y_position + 37)], fill=(0, 0, 0))

                    draw.text((170, y_position), line, fill=(0, 0, 0), font=point_font)
                else:

                    draw.text((170, y_position), line, fill=(0, 0, 0), font=point_font)
                
                y_position += 220
            
            y_position += 160  
    

    img.save(output_path)

if __name__ == '__main__':
    if not os.path.exists('static'):
        os.makedirs('static')
    

    init_db()
    
    app.run(debug=True, port=8080)